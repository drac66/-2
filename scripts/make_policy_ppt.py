from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TITLE_COLOR = RGBColor(31, 78, 121)
BODY_COLOR = RGBColor(51, 51, 51)
ACCENT = RGBColor(0, 112, 192)

slides = [
    ("夯基固本启新程，数字赋能向未来", [
        "乘势而上，接续推进中国式现代化与数字中国建设",
        "班级｜第X组｜成员｜日期",
        "动画：本页元素统一‘飞入-自底部’"
    ]),
    ("为什么要把两个主题合讲？", [
        "中国式现代化回答‘发展走向哪里’",
        "数字中国回答‘发展如何提速提质’",
        "两者关系：目标牵引路径，路径反哺目标",
        "一句话：现代化定方向，数字化强引擎"
    ]),
    ("中国式现代化：方向与价值", [
        "人口规模巨大的现代化",
        "全体人民共同富裕",
        "物质文明与精神文明相协调",
        "人与自然和谐共生",
        "走和平发展道路"
    ]),
    ("夯实基础：推进现代化的底盘", [
        "制度基础：治理效能与执行能力",
        "产业基础：现代化产业体系与实体经济韧性",
        "民生基础：教育、医疗、就业、社保托底",
        "安全基础：粮食、能源、产业链、数据安全"
    ]),
    ("全面发力：高质量发展的四个抓手", [
        "创新：科技自立自强，发展新质生产力",
        "协调：区域协同、城乡融合、产业联动",
        "绿色：低碳转型、生态价值转化",
        "开放：更高水平开放与规则对接"
    ]),
    ("数字中国：为什么是‘乘数效应’？", [
        "数字化同时作用于治理、生产、生活三大系统",
        "重构流程、重组资源、重塑效率",
        "数字中国效能=基础设施×数据要素×技术创新×场景落地"
    ]),
    ("建设路径：一底座、两能力、三场景", [
        "一底座：算力网络、数据平台、数字基础设施",
        "两能力：数据治理能力 + 技术创新能力",
        "三场景：数字政府、数字经济、数字社会"
    ]),
    ("场景一：数字政府（治理提效）", [
        "政务流程线上化、标准化、协同化",
        "群众办事少跑腿、少材料、少等待",
        "数据辅助决策，提升公共服务精准度"
    ]),
    ("场景二：数字经济（产业升级）", [
        "制造业数字化改造：可视化、柔性化、智能化",
        "传统产业与平台经济深度融合",
        "中小企业上云用数赋智，提升抗风险能力"
    ]),
    ("场景三：数字社会（普惠便民）", [
        "在线教育、互联网医疗、智慧交通持续优化",
        "公共服务覆盖更广，城乡可及性提升",
        "重视数字鸿沟治理，确保重点群体不掉队"
    ]),
    ("现实挑战：发展与安全双平衡", [
        "核心技术与关键环节仍需突破",
        "数据流通与隐私保护边界要更清晰",
        "地区、行业、人群数字能力差异仍存在"
    ]),
    ("结语：把目标变成可执行行动", [
        "中国式现代化是长期目标，需夯基固本",
        "数字中国是关键引擎，正在持续赋能",
        "把质量、效率、公平、安全统一起来",
        "谢谢大家，欢迎批评指正"
    ]),
]

for i, (title, bullets) in enumerate(slides):
    s = prs.slides.add_slide(prs.slide_layouts[6])

    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(1.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(36 if i == 0 else 32)
    r.font.bold = True
    r.font.color.rgb = TITLE_COLOR

    ln = s.shapes.add_shape(1, Inches(0.7), Inches(1.45), Inches(11.8), Inches(0.05))
    ln.fill.solid()
    ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background()

    bb = s.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(11.3), Inches(4.9))
    btf = bb.text_frame
    btf.word_wrap = True
    btf.clear()

    for j, b in enumerate(bullets):
        pp = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
        pp.text = b
        pp.level = 0
        pp.font.size = Pt(24 if i == 0 else 22)
        pp.font.color.rgb = BODY_COLOR
        pp.space_after = Pt(12)

    ft = s.shapes.add_textbox(Inches(0.7), Inches(6.9), Inches(12), Inches(0.4))
    fp = ft.text_frame.paragraphs[0]
    fp.text = "统一动画建议：标题与要点依次飞入（自底部），持续0.5秒，按0.1秒递增延迟"
    fp.font.size = Pt(12)
    fp.font.color.rgb = RGBColor(100, 100, 100)
    fp.alignment = PP_ALIGN.RIGHT

out = r"C:\Users\damaster\.openclaw\workspace\形式与政策_讨论课_第二版_可直接播放.pptx"
prs.save(out)
print(out)
